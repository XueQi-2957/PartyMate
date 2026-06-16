"""
PartyMate 文件导出引擎

支持从内容方案生成 PPT，从会议纪要生成 Word 文档。
"""

from __future__ import annotations

import json
import os
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

try:
    from docx import Document
    from docx.shared import Pt as DocxPt, RGBColor as DocxRGB, Inches as DocxInches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


# ---------- 输出目录 ---------- #

OUTPUT_DIR = Path(os.getcwd()) / "output"


def _ensure_output_dir():
    """确保输出目录存在"""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def _timestamp() -> str:
    """生成时间戳"""
    return date.today().strftime("%Y%m%d")


# ====================================================================
#  PPT 导出 — 从内容方案生成演示文稿
# ====================================================================

# 配色方案 — 红金主色调（党务风格）
C_PRIMARY = RGBColor(0xC0, 0x1B, 0x28)     # 中国红
C_SECONDARY = RGBColor(0xD4, 0xA0, 0x2C)    # 金色
C_DARK = RGBColor(0x2C, 0x2C, 0x2C)         # 深灰
C_BODY = RGBColor(0x33, 0x33, 0x33)          # 正文灰
C_LIGHT = RGBColor(0xF5, 0xF0, 0xEB)         # 米白底
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_slide_bg(slide, color: RGBColor):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=C_DARK, alignment=PP_ALIGN.LEFT,
                  font_name="微软雅黑"):
    """在幻灯片上添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_frame(slide, left, top, width, height, items,
                      font_size=14, color=C_BODY):
    """添加带项目符号的文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = Pt(6)
    return tf


def _make_cover_slide(prs, topic: str, meeting_type: str):
    """封面页 — 红底金字"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, C_PRIMARY)

    # 顶部装饰线
    _add_text_box(slide, 1.0, 1.0, 8.0, 0.3,
                  "━" * 60, font_size=8, color=C_SECONDARY,
                  alignment=PP_ALIGN.CENTER)

    # 标题
    _add_text_box(slide, 1.0, 2.2, 8.0, 1.5,
                  topic, font_size=32, bold=True,
                  color=C_WHITE, alignment=PP_ALIGN.CENTER)

    # 副标题
    _add_text_box(slide, 1.0, 3.8, 8.0, 0.6,
                  f"——  {meeting_type} 学习材料  ——",
                  font_size=16, color=C_SECONDARY,
                  alignment=PP_ALIGN.CENTER)

    # 底部信息
    _add_text_box(slide, 1.0, 5.6, 8.0, 0.5,
                  f"PartyMate · {date.today().strftime('%Y年%m月%d日')}",
                  font_size=11, color=RGBColor(0xFF, 0xCC, 0xCC),
                  alignment=PP_ALIGN.CENTER)

    # 底部装饰线
    _add_text_box(slide, 1.0, 6.5, 8.0, 0.3,
                  "━" * 60, font_size=8, color=C_SECONDARY,
                  alignment=PP_ALIGN.CENTER)


def _make_section_slide(prs, title: str, bullet_items: list[str]):
    """内容页 — 白底深色字"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_LIGHT)

    # 顶部标题栏
    _add_text_box(slide, 0.3, 0.3, 9.4, 0.7,
                  title, font_size=24, bold=True,
                  color=C_PRIMARY)

    # 标题下分隔线
    _add_text_box(slide, 0.3, 1.0, 9.4, 0.05,
                  "─" * 80, font_size=4, color=C_SECONDARY)

    # 内容
    _add_bullet_frame(slide, 0.5, 1.3, 9.0, 4.5,
                      bullet_items, font_size=15, color=C_BODY)


def _make_closing_slide(prs, topic: str):
    """结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_PRIMARY)

    _add_text_box(slide, 1.0, 2.0, 8.0, 1.0,
                  "感谢聆听", font_size=36, bold=True,
                  color=C_WHITE, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, 1.0, 3.5, 8.0, 0.8,
                  f"「{topic}」学习到此结束",
                  font_size=16, color=C_SECONDARY,
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, 1.0, 5.5, 8.0, 0.5,
                  "PartyMate · 党务智能助手",
                  font_size=12, color=RGBColor(0xFF, 0xCC, 0xCC),
                  alignment=PP_ALIGN.CENTER)


def export_content_pptx(content_json: str, output_path: str | None = None) -> str:
    """
    将内容方案导出为 PPTX 文件

    Args:
        content_json: 内容方案的 JSON 字符串（来自 generate_meeting_content）
        output_path: 输出路径，None 则自动生成

    Returns:
        生成的 PPTX 文件路径
    """
    from pptx import Presentation

    try:
        data = json.loads(content_json)
    except json.JSONDecodeError:
        return f"❌ 无法解析内容方案 JSON：{content_json[:100]}"

    topic = data.get("topic", "党务学习")
    meeting_type = data.get("meeting_type", "党课")
    ppt_outline = data.get("ppt_outline", [])
    discussion_questions = data.get("discussion_questions", [])

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 封面
    _make_cover_slide(prs, topic, meeting_type)

    # 内容页面
    for slide_data in ppt_outline:
        title = slide_data.get("title", "")
        content_text = slide_data.get("content", "")
        lines = [l.strip() for l in content_text.split("\n") if l.strip()]
        # 去掉首行（已体现在标题中）
        if lines and "封面" in slide_data.get("type", ""):
            continue  # 封面已在前面生成
        _make_section_slide(prs, title, lines)

    # 讨论题页
    if discussion_questions:
        _make_section_slide(
            prs, "讨论环节",
            [f"Q{i+1}. {q}" for i, q in enumerate(discussion_questions)]
        )

    # 结束页
    _make_closing_slide(prs, topic)

    # 保存
    _ensure_output_dir()
    if output_path:
        out = Path(output_path)
    else:
        safe_topic = topic.replace(" ", "").replace("/", "、")[:30]
        out = OUTPUT_DIR / f"{safe_topic}_党课PPT_{_timestamp()}.pptx"

    prs.save(str(out))
    return str(out.resolve())


# ====================================================================
#  Word 导出 — 从会议纪要生成正式文档
# ====================================================================


def _docx_add_title(doc, text: str, level: int = 0):
    """添加标题（带字体设置）"""
    if level == 0:
        heading = doc.add_heading(text, level=0)
        for run in heading.runs:
            run.font.color.rgb = DocxRGB(0xC0, 0x1B, 0x28)
    else:
        heading = doc.add_heading(text, level=level)
        for run in heading.runs:
            run.font.color.rgb = DocxRGB(0x2C, 0x2C, 0x2C)


def _docx_add_para(doc, text: str, bold=False, size=11, color=None, align=None):
    """添加段落"""
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.size = DocxPt(size)
    run.font.name = "仿宋"
    run.bold = bold
    if color:
        run.font.color.rgb = DocxRGB(*color)
    if align:
        p.alignment = align
    return p


def export_meeting_docx(summary_json: str, output_path: str | None = None) -> str:
    """
    将会议纪要导出为 Word 文档（.docx）

    Args:
        summary_json: 会议纪要的 JSON 字符串（来自 parse_meeting_notes）
        output_path: 输出路径，None 则自动生成

    Returns:
        生成的 DOCX 文件路径
    """
    try:
        data = json.loads(summary_json)
    except json.JSONDecodeError:
        return f"❌ 无法解析会议纪要 JSON：{summary_json[:100]}"

    info = data.get("meeting_info", {})
    topics = data.get("topics", [])
    action_items = data.get("action_items", [])
    key_decisions = data.get("summary", {}).get("key_decisions", [])

    doc = Document()

    # === 标题 ===
    meeting_type = info.get("type", "会议")
    meeting_date = info.get("date", "____年__月__日")
    _docx_add_title(doc, f"{meeting_type} 会议纪要", level=0)

    # === 会议基本信息表 ===
    doc.add_paragraph()  # 空行
    info_items = [
        ("会议时间", meeting_date),
        ("会议地点", info.get("location", "") or "（请填写）"),
        ("主持人", info.get("chair", "") or "（请填写）"),
        ("记录人", info.get("recorder", "") or "（请填写）"),
        ("出席人员", ", ".join(info.get("attendees", [])) or "（请填写）"),
        ("会议类型", meeting_type),
    ]
    table = doc.add_table(rows=len(info_items) + 1, cols=2)
    table.style = "Table Grid"

    # 表头
    for i, h in enumerate(["项目", "内容"]):
        cell = table.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = DocxPt(10)

    # 数据行
    for row_idx, (label, value) in enumerate(info_items, 1):
        table.rows[row_idx].cells[0].text = label
        table.rows[row_idx].cells[1].text = value or "（未填写）"
        for cell in table.rows[row_idx].cells:
            for p in cell.paragraphs:
                for r in p.runs:
                    r.font.size = DocxPt(10)

    # === 会议议程 ===
    doc.add_paragraph()
    _docx_add_title(doc, "一、会议议程", level=1)
    if topics:
        for i, topic in enumerate(topics, 1):
            _docx_add_para(doc, f"{i}. {topic.get('title', '（议题）')}", bold=True)
            for point in topic.get("discussion_points", []):
                _docx_add_para(doc, f"    • {point}", size=10.5)
    else:
        _docx_add_para(doc, "（无议题信息）")

    # === 决议事项 ===
    doc.add_paragraph()
    _docx_add_title(doc, "二、决议事项", level=1)
    if key_decisions and key_decisions != ["（请AI补充决议事项总结）"]:
        for decision in key_decisions:
            _docx_add_para(doc, f"• {decision}", size=10.5)
    else:
        _docx_add_para(doc, "（请补充决议事项）")

    # === 待办事项 ===
    doc.add_paragraph()
    _docx_add_title(doc, "三、待办事项", level=1)
    if action_items and action_items[0].get("task", "").startswith("（未识别"):
        _docx_add_para(doc, "（请补充待办事项）")
    else:
        table2 = doc.add_table(rows=len(action_items) + 1, cols=4)
        table2.style = "Table Grid"
        for i, h in enumerate(["序号", "待办事项", "负责人", "截止日期"]):
            cell = table2.rows[0].cells[i]
            cell.text = h
            for p in cell.paragraphs:
                for r in p.runs:
                    r.bold = True
                    r.font.size = DocxPt(10)
        for row_idx, item in enumerate(action_items, 1):
            table2.rows[row_idx].cells[0].text = str(row_idx)
            table2.rows[row_idx].cells[1].text = item.get("task", "")
            table2.rows[row_idx].cells[2].text = item.get("assignee", "")
            table2.rows[row_idx].cells[3].text = item.get("deadline", "")

    # === 落款 ===
    doc.add_paragraph()
    doc.add_paragraph()
    _docx_add_para(doc, f"{info.get('chair', '主持人')}（签字）", align=2)
    doc.add_paragraph()
    _docx_add_para(doc, date.today().strftime("%Y年%m月%d日"), align=2)

    # 保存
    _ensure_output_dir()
    if output_path:
        out = Path(output_path)
    else:
        safe_type = meeting_type.replace(" ", "")[:20]
        out = OUTPUT_DIR / f"{safe_type}会议纪要_{_timestamp()}.docx"

    doc.save(str(out))
    return str(out.resolve())


# ── 发展党员材料清单导出 ─────────────────────────────────────────


def export_member_material_checklist(
    member: dict, output_path: str | None = None
) -> str:
    """导出单个成员的发展党员材料清单 DOCX。

    Args:
        member: member dict from Repository.get_member() (includes 'materials' key)
        output_path: 可选指定输出路径

    Returns:
        str: 输出文件路径
    """
    doc = Document()

    # ── 标题 ──
    _docx_add_para(doc, "发展党员材料清单", size=18, bold=True, align=1)
    doc.add_paragraph()

    # ── 成员信息 ──
    info = doc.add_table(rows=5, cols=2, style="Table Grid")
    labels = ["姓名", "年级", "专业", "学号", "当前阶段"]
    keys = ["name", "grade", "major", "student_id", "stage"]
    for i, (label, key) in enumerate(zip(labels, keys)):
        val = member.get(key, "")
        if key == "stage":
            val = stage_labels.get(val, val)
        info.cell(i, 0).text = label
        info.cell(i, 1).text = str(val)

    doc.add_paragraph()

    # ── 阶段时间 ──
    date_fields = [
        ("apply_date", "递交申请日期"),
        ("activist_date", "确定积极分子日期"),
        ("candidate_date", "确定发展对象日期"),
        ("probationary_date", "接收预备党员日期"),
        ("full_member_date", "转正日期"),
    ]
    for field, label in date_fields:
        val = member.get(field, "")
        if val:
            _docx_add_para(doc, f"{label}：{val}")

    doc.add_paragraph()

    # ── 材料清单表格 ──
    materials = member.get("materials", [])
    if not materials:
        _docx_add_para(doc, "暂无材料记录")
    else:
        stage_order = ["applicant", "activist", "candidate", "probationary", "full_member"]
        # Group materials by stage
        grouped = {s: [] for s in stage_order}
        for mat in materials:
            s = mat.get("stage", "")
            if s in grouped:
                grouped[s].append(mat)

        table = doc.add_table(rows=1, cols=5, style="Table Grid")
        # Header
        hdr = table.rows[0].cells
        for i, txt in enumerate(["序号", "材料名称", "所属阶段", "状态", "提交日期"]):
            hdr[i].text = txt
            for p in hdr[i].paragraphs:
                for r in p.runs:
                    r.bold = True

        idx = 0
        submitted = 0
        total_required = 0
        for stage_key in stage_order:
            mats = grouped.get(stage_key, [])
            if not mats:
                continue
            # Stage header row
            row = table.add_row()
            stage_label = stage_labels.get(stage_key, stage_key)
            row.cells[0].merge(row.cells[4])
            row.cells[0].text = f"【{stage_label}】"
            for p in row.cells[0].paragraphs:
                for r in p.runs:
                    r.bold = True

            for mat in mats:
                idx += 1
                is_req = mat.get("is_required", 1)
                is_sub = mat.get("is_submitted", 0)
                if is_req:
                    total_required += 1
                if is_sub:
                    submitted += 1
                row = table.add_row()
                row.cells[0].text = str(idx)
                row.cells[1].text = mat.get("material_name", "")
                row.cells[2].text = stage_label
                row.cells[3].text = "✅ 已提交" if is_sub else "⬜ 未提交"
                row.cells[4].text = mat.get("submitted_date", "")

        doc.add_paragraph()
        _docx_add_para(
            doc,
            f"📊 汇总：已提交 {submitted}/{total_required} 项必需材料",
            size=12, bold=True,
        )

    # ── 导出 ──
    _ensure_output_dir()
    if output_path:
        out = Path(output_path)
    else:
        name = member.get("name", "未知")
        out = OUTPUT_DIR / f"{name}_发展党员材料清单_{_timestamp()}.docx"

    doc.save(str(out))
    return str(out.resolve())


def export_member_full_report(member: dict, output_path: str | None = None) -> str:
    """导出完整发展党员报告（材料清单 + 时间线 + 进度）。

    Args:
        member: member dict with 'materials', 'events' keys
        output_path: 可选指定输出路径

    Returns:
        str: 输出文件路径
    """
    # First, export the material checklist as base
    doc = Document()

    _docx_add_para(doc, "发展党员全流程报告", size=18, bold=True, align=1)
    doc.add_paragraph()

    name = member.get("name", "未知")
    _docx_add_para(doc, f"姓名：{name}")
    _docx_add_para(doc, f"当前阶段：{stage_labels.get(member.get('stage',''), member.get('stage',''))}")
    doc.add_paragraph()

    # ── 时间线 ──
    _docx_add_para(doc, "一、时间线", size=14, bold=True)
    events = member.get("events", [])
    if events:
        for ev in sorted(events, key=lambda x: x.get("event_date", "")):
            status_icon = "✅" if ev.get("status") == "completed" else "📌"
            _docx_add_para(
                doc,
                f"  {status_icon} {ev['event_type']}  ({ev.get('event_date','')})",
            )
    else:
        _docx_add_para(doc, "  暂无事件记录")

    doc.add_paragraph()

    # ── 材料清单 ──
    _docx_add_para(doc, "二、材料清单", size=14, bold=True)
    materials = member.get("materials", [])
    if materials:
        stage_order = ["applicant", "activist", "candidate", "probationary", "full_member"]
        grouped = {s: [] for s in stage_order}
        for mat in materials:
            s = mat.get("stage", "")
            if s in grouped:
                grouped[s].append(mat)

        table = doc.add_table(rows=1, cols=5, style="Table Grid")
        hdr = table.rows[0].cells
        for i, txt in enumerate(["序号", "材料名称", "所属阶段", "状态", "提交日期"]):
            hdr[i].text = txt
            for p in hdr[i].paragraphs:
                for r in p.runs:
                    r.bold = True

        idx = 0
        submitted = 0
        total = 0
        for stage_key in stage_order:
            mats = grouped.get(stage_key, [])
            if not mats:
                continue
            row = table.add_row()
            stage_label = stage_labels.get(stage_key, stage_key)
            row.cells[0].merge(row.cells[4])
            row.cells[0].text = f"【{stage_label}】"
            for p in row.cells[0].paragraphs:
                for r in p.runs:
                    r.bold = True
            for mat in mats:
                idx += 1
                if mat.get("is_required"):
                    total += 1
                if mat.get("is_submitted"):
                    submitted += 1
                row = table.add_row()
                row.cells[0].text = str(idx)
                row.cells[1].text = mat.get("material_name", "")
                row.cells[2].text = stage_label
                row.cells[3].text = "✅" if mat.get("is_submitted") else "⬜"
                row.cells[4].text = mat.get("submitted_date", "")

        doc.add_paragraph()
        # Progress bar (text-based)
        pct = round(submitted / total * 100) if total > 0 else 0
        bar_len = 20
        filled = round(bar_len * pct / 100)
        bar = "█" * filled + "░" * (bar_len - filled)
        _docx_add_para(doc, f"📊 总体进度：{submitted}/{total} ({pct}%)")
        _docx_add_para(doc, f"  {bar}  {pct}%")
    else:
        _docx_add_para(doc, "  暂无材料记录")

    # ── 导出 ──
    _ensure_output_dir()
    if output_path:
        out = Path(output_path)
    else:
        out = OUTPUT_DIR / f"{name}_发展党员全流程报告_{_timestamp()}.docx"

    doc.save(str(out))
    return str(out.resolve())


# Ensure stage_labels is available for both functions above
stage_labels = {
    "applicant": "申请入党", "activist": "入党积极分子",
    "candidate": "发展对象", "probationary": "预备党员",
    "full_member": "正式党员",
}
